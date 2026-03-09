#!/usr/bin/env python3
"""Generate Wynisco AI Matching Service — PowerPoint Presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import datetime
import diagram_generator
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Colors ──
DARK = RGBColor(0x1E, 0x29, 0x3B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
GREEN = RGBColor(0x10, 0xB9, 0x81)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
GRAY = RGBColor(0x6B, 0x72, 0x80)
LIGHT_GRAY = RGBColor(0xF1, 0xF5, 0xF9)
RED = RGBColor(0xEF, 0x44, 0x44)

# Pre-generate diagrams
diagrams = diagram_generator.generate_all()


def _set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=DARK, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def _add_bullet_frame(slide, left, top, width, height, items, font_size=16,
                      color=DARK, bold_prefix=False):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        if isinstance(item, tuple):
            run_bold = p.add_run()
            run_bold.text = item[0] + " "
            run_bold.font.bold = True
            run_bold.font.size = Pt(font_size)
            run_bold.font.color.rgb = color
            run_bold.font.name = "Calibri"
            run_norm = p.add_run()
            run_norm.text = item[1]
            run_norm.font.size = Pt(font_size)
            run_norm.font.color.rgb = color
            run_norm.font.name = "Calibri"
        else:
            p.text = f"  {item}"
    return txBox


def _add_table(slide, left, top, width, height, headers, rows, header_color=BLUE):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(
        n_rows, n_cols, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    table = table_shape.table

    # Header row
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = "Calibri"
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color

    # Data rows
    for r_idx, row_data in enumerate(rows):
        for c_idx, val in enumerate(row_data):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = DARK
                p.font.name = "Calibri"
            if r_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY

    return table_shape


# ══════════════════════════════════════════════════════════
# SLIDE 1: Title
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
_set_slide_bg(slide, DARK)

_add_textbox(slide, 1, 1.8, 11, 1.2, "AI Resume Matching Service",
             font_size=44, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
_add_textbox(slide, 1, 3.2, 11, 0.8,
             "Integration, Database Schema & Deployment Guide",
             font_size=22, color=RGBColor(0x94, 0xA3, 0xB8), alignment=PP_ALIGN.CENTER)
_add_textbox(slide, 1, 5.0, 11, 0.5,
             f"Version 1.0  |  {datetime.date.today().strftime('%B %d, %Y')}",
             font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)
_add_textbox(slide, 1, 5.5, 11, 0.5, "Confidential — Internal Use Only",
             font_size=12, color=RED, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════
# SLIDE 2: The Problem
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_add_textbox(slide, 0.5, 0.3, 12, 0.7, "The Problem",
             font_size=32, bold=True, color=DARK)

_add_bullet_frame(slide, 0.8, 1.2, 5.5, 5.5, [
    ("Today:", "consultants manually review job postings, compare requirements against "
     "each candidate's resume, and make subjective judgments about fit"),
    ("It's slow:", "hours per candidate to review and compare against available jobs"),
    ("It's inconsistent:", "different consultants evaluate the same candidate-job pair differently"),
    ("It doesn't scale:", "as the number of candidates and jobs grows, manual review becomes a bottleneck"),
], font_size=15)

_add_textbox(slide, 7.0, 1.3, 5.5, 0.5, "The Solution",
             font_size=22, bold=True, color=GREEN)
_add_bullet_frame(slide, 7.0, 2.0, 5.5, 5.0, [
    ("New microservice", "that automatically scores every candidate against every job"),
    ("Reads existing data:", "resumes and job descriptions from the same PostgreSQL database"),
    ("Uses AI:", "converts text into numerical fingerprints, then uses an LLM to produce detailed analysis"),
    ("Produces:", "ranked match results with plain-English explanations — strengths, gaps, recommendation"),
    ("Does not replace anything:", "runs alongside SMO, can be started/stopped independently"),
], font_size=15)


# ══════════════════════════════════════════════════════════
# SLIDE 3: Service At a Glance
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_add_textbox(slide, 0.5, 0.3, 12, 0.7, "Service at a Glance",
             font_size=32, bold=True, color=DARK)

_add_table(slide, 0.5, 1.2, 12.3, 5.0,
    ["Aspect", "Detail"],
    [
        ["Service Name", "auto-matching-service"],
        ["Technology", "Python 3.11, FastAPI (async), SQLAlchemy 2.0"],
        ["Ports", "8002 (backend API), 3001 (dashboard UI)"],
        ["AI Models", "OpenAI text-embedding-3-small (converts text into 1,536-number fingerprints)\n"
         "Cerebras gpt-oss-120b (reads both sides and produces scoring + analysis)"],
        ["Database", "Same PostgreSQL instance as SMO (smo_local).\n"
         "Reads 6 existing tables, writes to 3 new tables it owns."],
        ["Caching", "Two-layer: Redis (individual AI results, 72hr, optional) + DB (full results, 24hr).\n"
         "Works without Redis — falls back to DB-only caching automatically."],
        ["Dashboard", "Next.js web application for browsing match results, viewing AI explanations\n"
         "(WhyCards), tuning score weights, and listing candidates/jobs."],
    ],
    header_color=BLUE,
)


# ══════════════════════════════════════════════════════════
# SLIDE 4: Architecture
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_add_textbox(slide, 0.5, 0.3, 12, 0.7, "How It Fits Into the Existing Architecture",
             font_size=32, bold=True, color=DARK)

_add_textbox(slide, 0.5, 0.9, 12, 0.5,
             "Two new components (Matching Backend + Dashboard) plug into the existing infrastructure by sharing "
             "the same PostgreSQL database. No changes to SMO code required.",
             font_size=14, color=GRAY)

slide.shapes.add_picture(
    diagrams["architecture"], Inches(1.5), Inches(1.5), Inches(10), Inches(5.5)
)


# ══════════════════════════════════════════════════════════
# SLIDE 5: What Connects to What
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_add_textbox(slide, 0.5, 0.3, 12, 0.7, "How Each Connection Works",
             font_size=32, bold=True, color=DARK)

_add_textbox(slide, 0.5, 0.9, 12, 0.5,
             "Every line in the architecture diagram explained — what data moves, in which direction, and why.",
             font_size=14, color=GRAY)

_add_table(slide, 0.5, 1.5, 12.3, 5.5,
    ["Connection", "Direction", "What Happens"],
    [
        ["Backend → PostgreSQL", "Read existing",
         "Reads candidate profiles (resume, skills, experience, salary), job postings "
         "(descriptions, requirements, locations), and employer info (names, blacklist flags). "
         "This is read-only — the matching service never modifies SMO-owned data."],
        ["Backend → PostgreSQL", "Write new",
         "Writes to 3 owned tables: candidate_embeddings (AI fingerprints of resumes), "
         "job_embeddings (AI fingerprints of jobs), match_scores_v2 (scored results with explanations). "
         "These tables are separate from SMO tables."],
        ["Backend → OpenAI", "Outbound HTTPS",
         "Sends resume or job description text, receives a vector of 1,536 numbers representing "
         "the meaning of that text. One-time per document — result is stored and reused."],
        ["Backend → Cerebras", "Outbound HTTPS",
         "Sends candidate profile + job details to AI. Returns structured analysis: score (0-100), "
         "strengths, skill gaps, experience fit, salary fit, recommendation. This becomes the WhyCard."],
        ["Backend → Redis", "Cache read/write",
         "Checks for cached LLM result before calling Cerebras. Saves ~2-3 seconds and API cost per hit. "
         "Optional — if not configured, the system works using database-only caching."],
        ["Dashboard → Backend", "REST API",
         "Next.js proxy rewrites /api/* requests to backend (port 8002). Frontend doesn't need "
         "to know the backend's address directly."],
        ["SMO ↔ Matching", "No direct connection",
         "The two backends never call each other's APIs. They are completely decoupled services "
         "that happen to share a database. Either can be updated or restarted independently."],
    ],
    header_color=GREEN,
)


# ══════════════════════════════════════════════════════════
# SLIDE 6: Existing Tables Read
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_add_textbox(slide, 0.5, 0.3, 12, 0.7, "Existing Tables the Service Reads (Read-Only)",
             font_size=32, bold=True, color=DARK)

_add_textbox(slide, 0.5, 0.9, 12, 0.5,
             "Rather than duplicating data, the matching service reads directly from SMO's tables using "
             "read-only mappings. It always sees the latest data — when a consultant updates a candidate or a new job is scraped, "
             "the matching service picks up the changes automatically.",
             font_size=13, color=GRAY)

_add_table(slide, 0.5, 1.7, 12.3, 5.0,
    ["Table", "What It Reads", "Why It Needs It"],
    [
        ["users", "id (UUID), email, first_name, last_name, role, pod_id",
         "Identifies which users are candidates (filters by role). The id is the primary key "
         "that links everything — embeddings and match scores both reference this id."],
        ["candidate_profiles", "resume_plain_text, skills[], job_titles[],\n"
         "years_of_experience, location, expected_salary_min/max",
         "Core input for matching. Resume text goes to OpenAI for embedding. Skills, experience, "
         "location, salary used in hard filters and sent to LLM for detailed analysis."],
        ["jobs", "id, job_title, job_description, location,\n"
         "salary, experience, created_at, employer_id",
         "The other side of the match. Description goes to OpenAI for embedding. Title, location, "
         "salary, experience used in filtering and LLM analysis. created_at filters out old posts."],
        ["employers", "id, name, do_not_apply, h1b_filed,\n"
         "everified, location",
         "Every job links to an employer. do_not_apply is critical — if blacklisted, all their "
         "jobs are excluded from matching. Company name shown in match results."],
        ["pods", "id (UUID), pod_name",
         "Cohorts/batches of candidates. Dashboard lets you filter candidates by pod."],
        ["universities", "id, school_name, city, state",
         "Some candidates have a linked university. Shown in match results for background context."],
    ],
    header_color=PURPLE,
)


# ══════════════════════════════════════════════════════════
# SLIDE 7: New Tables Created
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_add_textbox(slide, 0.5, 0.3, 12, 0.7, "3 New Tables the Service Creates",
             font_size=32, bold=True, color=DARK)

_add_textbox(slide, 0.5, 0.9, 12, 0.5,
             "These tables store the AI-generated data. They live in the same PostgreSQL database as SMO but are "
             "managed by a separate migration chain — they never interfere with SMO's tables.",
             font_size=13, color=GRAY)

_add_table(slide, 0.5, 1.5, 12.3, 5.5,
    ["Table", "Purpose", "Key Columns", "Links To"],
    [
        ["candidate_embeddings",
         "AI fingerprint of each resume.\n"
         "Resume text is sent to OpenAI which returns\n"
         "1,536 numbers capturing its meaning.\n"
         "Two similar resumes have similar fingerprints.",
         "candidate_id (UUID, unique)\nembedding_vector (Vector 1536)\n"
         "source_text_hash (skip if unchanged)\nmodel_name, created_at, updated_at",
         "candidate_id\n→ users.id\n(one-to-one)"],
        ["job_embeddings",
         "AI fingerprint of each job description.\n"
         "Same concept as candidate_embeddings.\n"
         "During matching, candidate and job\n"
         "fingerprints are compared for similarity.",
         "job_id (Integer, unique)\nembedding_vector (Vector 1536)\n"
         "source_text_hash\nmodel_name, created_at, updated_at",
         "job_id\n→ jobs.id\n(one-to-one)"],
        ["match_scores_v2",
         "Final match results + 24hr cache.\n"
         "One row per candidate-job pair per direction.\n"
         "Contains scores AND the WhyCard —\n"
         "the AI's plain-English explanation.",
         "candidate_id (UUID), job_id (Integer)\n"
         "match_type ('candidate_to_jobs' / 'job_to_candidates')\n"
         "semantic_score, llm_score, combined_score\n"
         "score_breakdown (JSONB — the WhyCard)\n"
         "expires_at (created_at + 24 hours)",
         "candidate_id\n→ users.id\njob_id\n→ jobs.id\n(many-to-one)"],
    ],
    header_color=ORANGE,
)


# ══════════════════════════════════════════════════════════
# SLIDE 8: How New Tables Connect
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_add_textbox(slide, 0.5, 0.3, 12, 0.7, "How New Tables Connect to Existing Tables",
             font_size=32, bold=True, color=DARK)

_add_textbox(slide, 0.5, 0.9, 12, 0.4,
             "Each new table references SMO data through candidate_id and job_id columns.",
             font_size=14, color=GRAY)

_add_textbox(slide, 0.5, 1.4, 4, 0.4, "candidate_embeddings",
             font_size=18, bold=True, color=BLUE)
_add_textbox(slide, 0.5, 1.8, 5.5, 0.7,
             "candidate_id (UUID) → users.id\n"
             "Reads the candidate's profile from users + candidate_profiles, "
             "generates an embedding from their resume, stores result here. "
             "One row per candidate.",
             font_size=12, color=DARK)

_add_textbox(slide, 0.5, 2.6, 4, 0.4, "job_embeddings",
             font_size=18, bold=True, color=GREEN)
_add_textbox(slide, 0.5, 3.0, 5.5, 0.7,
             "job_id (Integer) → jobs.id\n"
             "Reads the job description from jobs (which links to employers "
             "via jobs.employer_id → employers.id), generates an embedding, "
             "stores result here. One row per job.",
             font_size=12, color=DARK)

_add_textbox(slide, 0.5, 3.8, 4, 0.4, "match_scores_v2",
             font_size=18, bold=True, color=ORANGE)
_add_textbox(slide, 0.5, 4.2, 5.5, 0.9,
             "candidate_id (UUID) → users.id\n"
             "job_id (Integer) → jobs.id\n"
             "Connects a candidate to a job with scores and the AI explanation "
             "(WhyCard). One candidate can have many rows (one per job matched). "
             "100 candidates × 200 jobs = up to 20,000 rows.",
             font_size=12, color=DARK)

slide.shapes.add_picture(
    diagrams["db_ownership"], Inches(6.5), Inches(1.3), Inches(6.3), Inches(3.8)
)

_add_textbox(slide, 0.5, 5.3, 12, 0.3, "The Complete Reference Chain",
             font_size=16, bold=True, color=DARK)

_add_table(slide, 0.5, 5.7, 12.3, 1.6,
    ["Existing Table (SMO, hard FK)", "→", "New Table (Matching, logical reference)"],
    [
        ["users.id ← candidate_profiles.user_id", "⇠", "candidate_embeddings.candidate_id  &  match_scores_v2.candidate_id"],
        ["jobs.id → jobs.employer_id → employers.id", "⇠", "job_embeddings.job_id  &  match_scores_v2.job_id"],
    ],
    header_color=PURPLE,
)


# ══════════════════════════════════════════════════════════
# SLIDE 9: Why No Hard Foreign Keys?
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_add_textbox(slide, 0.5, 0.3, 12, 0.7, "Why No Hard Foreign Key Constraints?",
             font_size=32, bold=True, color=DARK)

_add_textbox(slide, 0.5, 1.0, 12, 0.6,
             "The new tables reference users.id and jobs.id but intentionally don't declare "
             "ForeignKey() constraints in the database. This is a deliberate design choice, not an oversight.",
             font_size=16, color=DARK)

_add_bullet_frame(slide, 0.8, 1.8, 11, 5.0, [
    ("Separate migration chains —",
     "SMO tracks its schema version in the alembic_version table; the matching service tracks "
     "its version in alembic_version_auto_matcher. If the new tables declared hard FK constraints "
     "pointing at SMO tables, Alembic would try to manage tables it doesn't own, causing migration conflicts."),
    ("Decoupled deployments —",
     "Because there are no hard FK constraints, either service can be deployed, rolled back, or "
     "restarted independently. You can drop and recreate all 3 matching tables without affecting "
     "any SMO data, and you can update SMO's schema without worrying about breaking the matching "
     "service's references."),
    ("Read-only contract —",
     "The matching service uses read-only SQLAlchemy mappings (defined in app/models/readonly.py "
     "with extend_existing=True) for all SMO tables. It physically cannot write to or modify SMO "
     "data, even accidentally. The relationship is enforced at the application layer (Python code), "
     "not at the database layer."),
], font_size=15)

_add_textbox(slide, 0.8, 5.7, 11, 0.8,
             "In practice: if a candidate is deleted in SMO, their embedding row becomes orphaned — "
             "but this is harmless. The matching service simply ignores orphaned rows. No errors, "
             "no cascading failures, no manual cleanup needed.",
             font_size=14, color=GRAY)


# ══════════════════════════════════════════════════════════
# SLIDE 10: Step-by-Step: How the Dashboard Reads a Match
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_add_textbox(slide, 0.5, 0.3, 12, 0.7,
             "How the Dashboard Assembles a Match Result",
             font_size=32, bold=True, color=DARK)

_add_textbox(slide, 0.5, 0.9, 12, 0.5,
             'When you see "John Doe matched to Software Engineer at Acme Corp, score 85/100" on the dashboard, '
             "here's the 5-step chain the system follows:",
             font_size=14, color=GRAY)

_add_bullet_frame(slide, 0.8, 1.5, 5.5, 5.5, [
    ("Step 1:", "Read the match score row from match_scores_v2\n→ gets the combined score and WhyCard (AI explanation)"),
    ("Step 2:", "Follow candidate_id → users.id\n→ gets the candidate's name and email"),
    ("Step 3:", "Follow users.id → candidate_profiles.user_id\n→ gets resume text, skills, experience, location, salary"),
    ("Step 4:", "Follow job_id → jobs.id\n→ gets the job title and description"),
    ("Step 5:", "Follow jobs.employer_id → employers.id\n→ gets the company name (\"Acme Corp\")"),
], font_size=14)

slide.shapes.add_picture(
    diagrams["match_chain"], Inches(6.5), Inches(1.2), Inches(6.5), Inches(5.5)
)


# ══════════════════════════════════════════════════════════
# SLIDE 11: AI Pipeline
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_add_textbox(slide, 0.5, 0.3, 12, 0.7, "AI Matching Pipeline — Three Layers",
             font_size=32, bold=True, color=DARK)

_add_textbox(slide, 0.5, 0.9, 12, 0.5,
             "The pipeline works as a funnel: start with all jobs, progressively filter and rank, "
             "produce a short list with detailed explanations. Each layer is more expensive but more intelligent.",
             font_size=14, color=GRAY)

_add_table(slide, 0.5, 1.5, 12.3, 4.5,
    ["Layer", "What It Does", "Data Used", "Output"],
    [
        ["1. Hard Filter\n(SQL, ~100ms)",
         "Eliminates obviously bad matches using\n"
         "simple rules. A job requiring 5 years of\n"
         "experience won't match a candidate with 2.\n"
         "Blacklisted employers excluded entirely.",
         "candidate_profiles\n(experience, location, salary)\n"
         "jobs (requirements)\n"
         "employers (do_not_apply)",
         "Up to 8,000 viable jobs\n(from tens of thousands)"],
        ["2. Semantic Search\n(pgvector, ~200ms)",
         "Compares meaning of resume vs job\n"
         "descriptions using pre-computed embeddings.\n"
         "Finds non-obvious matches that keyword\n"
         "search misses (e.g. 'Python dev' → 'Django engineer').",
         "candidate_embeddings\n.embedding_vector\n\n"
         "job_embeddings\n.embedding_vector",
         "Top 50 most similar jobs\nwith semantic score (0-1)"],
        ["3. LLM Analysis\n(Cerebras, ~2-3s each)",
         "AI reads both the resume and job description\n"
         "in detail. Produces structured analysis:\n"
         "score, strengths, skill gaps, experience fit,\n"
         "salary fit, concerns, recommendation.",
         "candidate name, skills,\nexperience, location, salary\n\n"
         "job title, description,\nlocation, salary, experience",
         "Top 10 scored matches\neach with a WhyCard\n(plain-English explanation)"],
    ],
    header_color=ORANGE,
)

_add_textbox(slide, 0.5, 6.2, 12, 0.8,
             "Final score = (semantic_score × 50%) + (llm_score × 50%)  —  weights are tunable from the dashboard in real-time",
             font_size=16, bold=True, color=DARK, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════
# SLIDE 12: Caching
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_add_textbox(slide, 0.5, 0.3, 12, 0.7, "Two-Layer Caching",
             font_size=32, bold=True, color=DARK)

_add_textbox(slide, 0.5, 0.9, 12, 0.4,
             "LLM analysis is the slowest/most expensive step (~2-3 seconds per call). "
             "Without caching, matching 100 candidates × 200 jobs = 20,000 LLM calls. Caching eliminates repeat work.",
             font_size=13, color=GRAY)

slide.shapes.add_picture(
    diagrams["cache_flow"], Inches(0.2), Inches(1.4), Inches(6.5), Inches(5.3)
)

_add_bullet_frame(slide, 7.0, 1.5, 5.8, 3.5, [
    ("Layer 1 — DB Cache:", "Stores complete match results in match_scores_v2.\n"
     "24hr TTL. If the same match is requested again\n"
     "within 24 hours, the entire pipeline is skipped\n"
     "and cached result returned in ~5ms."),
    ("Layer 2 — Redis Cache:", "Stores individual LLM results per (candidate, job) pair.\n"
     "72hr TTL. If the same pair appears in a different\n"
     "request, the Cerebras API call is skipped (~1ms).\n"
     'E.g. "jobs for John" cached pair reused in\n'
     '"candidates for Acme SWE".'),
    ("Redis is optional:", "If REDIS_URL is not configured or Redis goes down,\n"
     "the service logs a warning and continues with\n"
     "DB-only caching. No errors, self-heals when Redis returns."),
], font_size=13)

_add_table(slide, 7.0, 5.3, 5.8, 1.8,
    ["Scenario", "Cost", "Response Time"],
    [
        ["First run: 100 × 200 jobs", "~$1-5", "2-3 min"],
        ["Repeat same batch", "$0", "~5ms"],
    ],
    header_color=RGBColor(0xDC, 0x26, 0x26),
)


# ══════════════════════════════════════════════════════════
# SLIDE 13: Database Ownership & Migration Safety
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_add_textbox(slide, 0.5, 0.3, 12, 0.7, "Database Ownership & Migration Safety",
             font_size=32, bold=True, color=DARK)

_add_textbox(slide, 0.5, 0.9, 12, 0.5,
             "Two services share one database — how do we make sure they don't break each other's tables?",
             font_size=16, color=DARK)

_add_table(slide, 0.5, 1.6, 12.3, 3.5,
    ["Aspect", "SMO Backend", "Matching Service"],
    [
        ["Migration tool", "Alembic", "Alembic"],
        ["Version table", "alembic_version", "alembic_version_auto_matcher"],
        ["Tables managed", "33 tables (all SMO data)",
         "3 tables only (candidate_embeddings,\njob_embeddings, match_scores_v2)"],
        ["Safety filter", "None needed (owns the whole database)",
         "OWNED_TABLES filter in env.py —\nautogenerate ignores all other tables"],
        ["Run independently?", "Yes — migrations never touch\nmatching tables",
         "Yes — migrations never touch\nSMO tables"],
    ],
    header_color=PURPLE,
)

_add_textbox(slide, 0.5, 5.5, 12, 1.2,
             "You can safely run 'alembic upgrade head' for both services in any order. "
             "They use separate version tables, so they track their migration history independently. "
             "Running one will never roll back or interfere with the other.\n\n"
             "The matching service's OWNED_TABLES filter explicitly restricts its migration autogeneration "
             "to only its 3 tables — even though it can see all 36 tables in the database, it will never "
             "generate scripts that touch SMO's tables.",
             font_size=14, color=GRAY)


# ══════════════════════════════════════════════════════════
# SLIDE 14: Deployment
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_add_textbox(slide, 0.5, 0.3, 12, 0.7, "Deployment",
             font_size=32, bold=True, color=DARK)

# Left: Local dev
_add_textbox(slide, 0.5, 1.2, 6, 0.5, "Local Development",
             font_size=20, bold=True, color=GREEN)
_add_bullet_frame(slide, 0.8, 1.8, 5.5, 4.0, [
    "cd auto-matching-service",
    "cp .env.example .env   (fill in API keys)",
    "pip install -r requirements.txt",
    "alembic upgrade head   (creates 3 new tables only)",
    "uvicorn app.main:app --reload --port 8002",
    "",
    "cd frontend && npm install && npm run dev",
    "(optional) Set REDIS_URL in .env for Redis caching",
], font_size=13, color=DARK)

# Right: Cloud Run
_add_textbox(slide, 6.8, 1.2, 6, 0.5, "Google Cloud Run",
             font_size=20, bold=True, color=BLUE)
_add_bullet_frame(slide, 7.1, 1.8, 5.5, 4.0, [
    "gcloud run deploy auto-matcher-backend",
    "gcloud run deploy auto-matcher-frontend",
    "",
    ("Required:", "Cloud SQL (PostgreSQL 15+ with pgvector extension)"),
    ("Optional:", "Memorystore Redis (Basic tier, 1GB — ~$7/mo)"),
    ("Secrets:", "DATABASE_URL, OPENAI_API_KEY, CEREBRAS_API_KEY"),
    ("Verify:", "Visit /docs on backend URL for API docs"),
], font_size=13, color=DARK)

# Bottom: Env vars
_add_textbox(slide, 0.5, 5.6, 12, 0.4, "Environment Variables (only 3 required — everything else has defaults)",
             font_size=16, bold=True, color=DARK)
_add_table(slide, 0.5, 6.0, 12.3, 1.3,
    ["Variable", "Required?", "Default", "Description"],
    [
        ["DATABASE_URL", "Yes", "—", "postgresql+asyncpg://user@host:5432/smo_local"],
        ["OPENAI_API_KEY", "Yes", "—", "For embedding generation (text-embedding-3-small)"],
        ["CEREBRAS_API_KEY", "Yes", "—", "For LLM scoring (gpt-oss-120b)"],
        ["REDIS_URL", "No", "(disabled)", "Leave empty — system works without it, falls back to DB-only cache"],
    ],
    header_color=DARK,
)


# ══════════════════════════════════════════════════════════
# SLIDE 15: Cost
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_add_textbox(slide, 0.5, 0.3, 12, 0.7, "Cost Estimates",
             font_size=32, bold=True, color=DARK)

_add_textbox(slide, 0.5, 1.2, 6, 0.5, "Infrastructure (Monthly)",
             font_size=20, bold=True, color=GREEN)

_add_table(slide, 0.5, 1.8, 5.8, 3.0,
    ["Service", "Cost", "Notes"],
    [
        ["Cloud Run (backend)", "$5–15/mo", "Scales to zero when idle"],
        ["Cloud Run (frontend)", "$3–8/mo", "Scales to zero when idle"],
        ["Cloud SQL (shared)", "$0 incremental", "Already running for SMO"],
        ["Redis (optional)", "$7–10/mo", "Can skip entirely"],
        ["Total", "$15–33/mo", ""],
    ],
    header_color=GREEN,
)

_add_textbox(slide, 6.8, 1.2, 6, 0.5, "API Costs (Per Matching Run)",
             font_size=20, bold=True, color=ORANGE)

_add_table(slide, 6.8, 1.8, 5.8, 3.0,
    ["Scenario", "Cost", "Time"],
    [
        ["10 candidates × 20 jobs", "~$0.21", "~30 seconds"],
        ["100 candidates × 200 jobs", "~$1-5", "~2-3 minutes"],
        ["Repeat (all cached)", "$0", "~5ms"],
    ],
    header_color=ORANGE,
)

_add_textbox(slide, 0.5, 5.3, 12, 1.2,
             "After the first run, repeat requests for the same candidate-job pairs "
             "cost $0 in API fees — results are served from cache (24hr DB cache + 72hr Redis cache).\n\n"
             "Total estimated cost: ~$15-33/month infrastructure + ~$1-5 per batch in API fees.\n"
             "The matching service uses the existing Cloud SQL database, so there is no incremental database cost.",
             font_size=15, color=GRAY, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════
output_path = "Wynisco_Matching_Service_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
